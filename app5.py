import streamlit as st
import base64
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv
from groq import Groq
from io import BytesIO
import tempfile
import zipfile
import xml.etree.ElementTree as ET
import time
from PIL import Image as Imagee
import sounddevice as sd
import numpy as np
import scipy.io.wavfile as wavfile
import speech_recognition as sr
from spire.presentation import Presentation
from spire.presentation.common import *
import wave
import struct

# Load environment variables
load_dotenv()

# Get API key from environment variable
api_key = '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'
if not api_key:
    raise ValueError("GROQ_API_KEY not found in environment variables")

client = Groq(api_key=api_key)

# Custom Formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    prompt = f"""Generate exactly 5 concise slide titles for a presentation on the topic: {topic}
    Rules:
    1. Provide only the titles, one per line
    2. Do not include any numbering or bullet points
    3. Each title should be brief and relevant to the topic
    4. Do not include any additional text in response  or explanations"""
    
    response = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content

    # Split the response into titles and filter out empty lines
    return [title.strip() for title in response_text.split("\n") if title.strip()][:5] 

def generate_slide_content(slide_title):
    prompt = f"""Generate exactly 6 bullet points for the slide titled: "{slide_title}"
    Rules:
    1. Each point must be a very short but crisp sentence
    2. Do not exceed 10 words per point
    3. Provide only the points, one per line
    4. Do not include any numbering or bullet point symbols
    5. Do not include any additional text from response or explanations"""
    
    response = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content

    # Split the response into points and filter out empty lines
    points = [point.strip() for point in response_text.split("\n") if point.strip()][:6]  # Ensure we get exactly 6 points
    
    # Join the points with newlines to create the slide content
    return "\n\n".join(points)

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    # Add the title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Add the "Contents" slide (second slide)
    contents_slide = prs.slides.add_slide(slide_layout)
    contents_slide.shapes.title.text = "Contents"

    # Add each slide title as bullet points on the Contents slide
    content_text = "\n\n".join(slide_titles)  # Join all slide titles into one text with line breaks
    contents_slide.shapes.placeholders[1].text = content_text

    # Customize font size for the "Contents" slide
    contents_slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    for paragraph in contents_slide.shapes.placeholders[1].text_frame.paragraphs:
        paragraph.font.size = SLIDE_FONT_SIZE

    # Add the rest of the slides with their respective content
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content
        
        # Customize font size for each slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
                    
                    
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[0])
    thank_you_slide.shapes.title.text = "Thank You"


    # Save the presentation
    os.makedirs('generated_ppt', exist_ok=True)
    ppt_path = os.path.join('generated_ppt', f'{topic}_presentation.pptx')
    prs.save(ppt_path)
    return ppt_path


def get_ppt_download_link(ppt_path):
    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()
    
    b64_ppt = base64.b64encode(ppt_contents).decode()

    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(ppt_path)}">Download the PowerPoint</a>'

def rate_ppt(ppt_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(ppt_file.getvalue())
        tmp_file_path = tmp_file.name

    try:
        with zipfile.ZipFile(tmp_file_path, 'r') as zip_ref:
            # Get total number of slides
            presentation_xml = zip_ref.read('ppt/presentation.xml')
            root = ET.fromstring(presentation_xml)
            total_slides = len(root.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'))

            # Initialize ratings
            slide_number_rating = 10
            bullet_point_rating = 10

            # Analyze each slide
            for i in range(1, total_slides + 1):
                try:
                    slide_xml = zip_ref.read(f'ppt/slides/slide{i}.xml')
                    slide_root = ET.fromstring(slide_xml)

                    # Check for slide number
                    slide_number = slide_root.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ph[@type="sldNum"]')
                    if slide_number is None:
                        slide_number_rating -= 0.5

                    # Count bullet points
                    bullet_points = slide_root.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    if len(bullet_points) > 7:
                        bullet_point_rating -= 1

                except KeyError:
                    # Slide doesn't exist, reduce rating
                    slide_number_rating -= 1

        # Ensure ratings don't go below 0
        slide_number_rating = max(0, slide_number_rating)
        bullet_point_rating = max(0, bullet_point_rating)

        # Calculate overall rating
        overall_rating = (slide_number_rating + bullet_point_rating) / 2

        return {
            "overall_rating": overall_rating,
            "slide_number_rating": slide_number_rating,
            "bullet_point_rating": bullet_point_rating,
            "total_slides": total_slides
        }

    finally:
        # Clean up the temporary file
        os.unlink(tmp_file_path)

# Function to convert PPT to images
def convert_ppt_to_images(ppt_file):
    # Save the uploaded file to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt_file:
        tmp_ppt_file.write(ppt_file.getvalue())  # Write the content of the uploaded file
        tmp_ppt_file_path = tmp_ppt_file.name    # Get the file path

    try:
        # Load the presentation from the saved file
        presentation = Presentation()
        presentation.LoadFromFile(tmp_ppt_file_path)

        images = []
        with tempfile.TemporaryDirectory() as tmpdirname:
            for i, slide in enumerate(presentation.Slides):
                fileName = f"slide_{i}.png"
                image = slide.SaveAsImageByWH(800, 450)
                file_path = os.path.join(tmpdirname, fileName)
                image.Save(file_path)
                image.Dispose()

                # Open the image, convert it to RGB (to ensure compatibility), and store it in memory
                with Imagee.open(file_path) as img:
                    images.append(img.copy().convert('RGB'))

        presentation.Dispose()
        return images

    finally:
        # Clean up the temporary file
        if os.path.exists(tmp_ppt_file_path):
            os.remove(tmp_ppt_file_path)


# Function to display slideshow
def display_slideshow(slides):
    if not slides:
        st.write("No slides to display.")
        return

    slideshow_placeholder = st.empty()
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        prev_button = st.button("Previous")
    with col2:
        next_button = st.button("Next")
    with col3:
        auto_play = st.checkbox("Auto-play")

    if 'slide_index' not in st.session_state:
        st.session_state.slide_index = 0

    if prev_button and st.session_state.slide_index > 0:
        st.session_state.slide_index -= 1
    elif next_button and st.session_state.slide_index < len(slides) - 1:
        st.session_state.slide_index += 1

    if auto_play:
        time.sleep(3)
        st.session_state.slide_index = (st.session_state.slide_index + 1) % len(slides)
        st.experimental_rerun()

    slideshow_placeholder.image(slides[st.session_state.slide_index], use_column_width=True)

# Function to record audio
def record_audio(duration, samplerate=16000):
    st.write("Recording...")
    recording = sd.rec(int(samplerate * duration), samplerate=samplerate, channels=1, dtype='float32')
    sd.wait()
    st.write("Recording complete.")
    return recording

# Function to save audio
def save_audio(recording, samplerate=16000):  # Changed to 16000 Hz
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
    
    with wave.open(temp_file.name, 'w') as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)  # 16-bit samples (2 bytes)
        wf.setframerate(samplerate)
        for sample in recording:
            wf.writeframes(struct.pack('<h', int(sample[0] * 32767)))
    
    return temp_file.name

# Function to transcribe audio using SpeechRecognition
def transcribe_audio(audio_file):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_file) as source:
        audio_data = recognizer.record(source)
        try:
            text = recognizer.recognize_google(audio_data)
            return text
        except sr.UnknownValueError:
            return "Audio is unclear. Could not transcribe."
        except sr.RequestError as e:
            return f"Could not request results; {e}"

# Main Streamlit App
st.title("AI-Powered Presentation Tool")

tab1, tab2, tab3 = st.tabs(["Generate Presentation", "Rate Presentation", "Presentation Training"])

with tab1:
    st.header("Generate AI-powered Presentation")
    topic = st.text_input("Enter the topic for your presentation")
    
    if st.button("Generate Presentation"):
        with st.spinner("Generating slide titles..."):
            slide_titles = generate_slide_titles(topic)

        slide_contents = []
        for title in slide_titles:
            with st.spinner(f"Generating content for slide: {title}"):
                slide_content = generate_slide_content(title)
                slide_contents.append(slide_content)

        ppt_path = create_presentation(topic, slide_titles, slide_contents)
        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(ppt_path), unsafe_allow_html=True)

with tab2:
    st.header("Rate Your Presentation")
    ppt_file = st.file_uploader("Upload your PowerPoint file", type=["pptx"])
    if ppt_file is not None:
        rating = rate_ppt(ppt_file)
        st.write(f"Total Slides: {rating['total_slides']}")
        st.write(f"Slide Number Rating: {rating['slide_number_rating']}/10")
        st.write(f"Bullet Point Rating: {rating['bullet_point_rating']}/10")
        st.write(f"Overall Rating: {rating['overall_rating']}/10")

with tab3:
    st.header("Slide Show with Audio Transcription")

    ppt_file = st.file_uploader("Upload your PowerPoint file for slideshow", type=["pptx"], key="slideshow_ppt")

    if ppt_file is not None:
        slides = convert_ppt_to_images(ppt_file)

        display_slideshow(slides)

        duration = st.number_input("Enter duration (in seconds) to record:", min_value=1, max_value=60)
        if st.button("Record Audio"):
            recording = record_audio(duration)
            st.write("Audio recorded successfully!")

            audio_file = save_audio(recording)
            st.write(f"Audio saved at: {audio_file}")

            transcription = transcribe_audio(audio_file)
            st.write("Transcription:")
            st.write(transcription)
